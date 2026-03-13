from typing import Any, Dict, List, Optional, Tuple, Union
from pathlib import Path
from pptx.presentation import Presentation
from pptx.util import Inches
from pptx.chart.chart import Chart
from pptx.shapes.picture import Picture
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image
import re

class FontProperties:
    """Classe pour représenter les propriétés de police."""
    
    def __init__(self, name: str = None, size: int = None,
                 bold: bool = None, italic: bool = None, underline: bool = None):
        self.name = name or "Calibri"
        self.size = size or 12
        self.bold = bold or False
        self.italic = italic or False
        self.underline = underline or False

class FormattedText:
    """Classe pour représenter du texte formaté."""
    
    def __init__(self, text: str, font_name: str = None, font_size: int = None,
                 bold: bool = None, italic: bool = None, underline: bool = None):
        self.text = text
        self.font = FontProperties(font_name, font_size, bold, italic, underline)

def format_text(text: str, font_size: Optional[int] = None, bullet_style: Optional[str] = None) -> str:
    """
    Format text for slide content.
    
    Args:
        text: Text to format
        font_size: Optional font size in points
        bullet_style: Optional bullet style
        
    Returns:
        Formatted text
    """
    formatted = text.strip()
    
    # Remove markdown-style formatting
    formatted = re.sub(r'\*\*(.*?)\*\*', r'\1', formatted)  # Bold
    formatted = re.sub(r'\*(.*?)\*', r'\1', formatted)      # Italic
    formatted = re.sub(r'__(.*?)__', r'\1', formatted)      # Underline
    
    if bullet_style and not formatted.startswith(bullet_style):
        formatted = f"{bullet_style} {formatted}"
    
    return formatted

def get_slide_dimensions(presentation: Presentation) -> Dict[str, int]:
    """Retourne les dimensions d'une diapositive.
    
    Args:
        presentation: La présentation à analyser
        
    Returns:
        Un dictionnaire contenant la largeur et la hauteur
    """
    width = int(presentation.slide_width)
    height = int(presentation.slide_height)
    return {"width": width, "height": height}

def get_slide_layouts(presentation: Presentation) -> List[str]:
    """Retourne la liste des mises en page disponibles.
    
    Args:
        presentation: La présentation à analyser
        
    Returns:
        Une liste des noms de mises en page
    """
    layouts = []
    for master in presentation.slide_masters:
        for layout in master.slide_layouts:
            if layout.name not in layouts:
                layouts.append(layout.name)
    return layouts

def get_slide_masters(presentation: Presentation) -> List[str]:
    """Retourne la liste des masques disponibles.
    
    Args:
        presentation: La présentation à analyser
        
    Returns:
        Une liste des noms de masques
    """
    return [master.name for master in presentation.slide_masters]

def get_slide_theme(presentation: Presentation) -> Optional[str]:
    """Retourne le thème actuel de la présentation.
    
    Args:
        presentation: La présentation à analyser
        
    Returns:
        Le nom du thème ou None si aucun thème n'est appliqué
    """
    if presentation.slide_masters:
        return presentation.slide_masters[0].name
    return None

def add_chart(
    slide: Any,
    data: Dict[str, Any],
    chart_type: Union[str, XL_CHART_TYPE] = "bar",
    position: Optional[Tuple[float, float, float, float]] = None,
    title: Optional[str] = None,
    style: Optional[int] = None
) -> Chart:
    """Ajoute un graphique à une slide - refactorisé pour réduire la complexité"""

    # Étape 1: Validation et conversion du type de graphique
    chart_type = _validate_chart_type(chart_type)

    # Étape 2: Préparation des données du graphique
    chart_data = _prepare_chart_data(data)

    # Étape 3: Configuration de la position
    if position is None:
        position = _setup_default_chart_position()

    # Étape 4: Création du graphique sur la slide
    chart = _create_chart_on_slide(slide, chart_type, position, chart_data)

    # Étape 5: Application du style et du titre
    _apply_chart_styling(chart, title, style)

    return chart

def _get_chart_types_mapping() -> Dict[str, XL_CHART_TYPE]:
    """Retourne le mapping des types de graphiques"""
    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }

def _validate_chart_type(chart_type: Union[str, XL_CHART_TYPE]) -> XL_CHART_TYPE:
    """Valide et convertit le type de graphique"""
    chart_types = _get_chart_types_mapping()

    if isinstance(chart_type, str):
        if chart_type.lower() not in chart_types:
            raise ValueError(f"Type de graphique invalide. Valeurs possibles : {list(chart_types.keys())}")
        return chart_types[chart_type.lower()]
    elif not isinstance(chart_type, XL_CHART_TYPE):
        raise ValueError("Le type de graphique doit être une chaîne ou un XL_CHART_TYPE")

    return chart_type

def _prepare_chart_data(data: Dict[str, Any]) -> ChartData:
    """Prépare les données pour le graphique"""
    chart_data = ChartData()
    chart_data.categories = data["labels"]
    chart_data.add_series("Série 1", data["values"])
    return chart_data

def _setup_default_chart_position() -> Tuple[float, float, float, float]:
    """Configure la position par défaut du graphique"""
    return (2, 2, 6, 4.5)  # x, y, width, height en pouces

def _create_chart_on_slide(slide, chart_type, position, chart_data) -> Chart:
    """Crée le graphique sur la slide"""
    x, y, cx, cy = position
    return slide.shapes.add_chart(
        chart_type,
        Inches(x), Inches(y),
        Inches(cx), Inches(cy),
        chart_data
    ).chart

def _apply_chart_styling(chart, title: Optional[str], style: Optional[int]):
    """Applique le style et le titre au graphique"""
    if title:
        chart.chart_title.text_frame.text = title

    if style is not None:
        chart.chart_style = style

def add_image(
    slide: Any,
    image_path: Union[str, Path],
    position: Optional[Tuple[float, float, float, float]] = None,
    maintain_aspect_ratio: bool = True,
    crop: Optional[Tuple[float, float, float, float]] = None,
    transparency: Optional[float] = None
) -> Picture:
    """Ajoute une image à une slide - refactorisé pour réduire la complexité"""
    image_path = Path(image_path)

    # Étape 1: Validation du chemin d'image
    _validate_image_path(image_path)

    # Étape 2: Calcul des dimensions et position
    if maintain_aspect_ratio:
        position = _calculate_image_dimensions_with_aspect_ratio(image_path, position)
    elif position is None:
        position = _setup_default_image_position()

    # Étape 3: Ajout de l'image à la slide
    image = _add_image_to_slide(slide, image_path, position)

    # Étape 4: Application des effets optionnels
    _apply_image_cropping(image, crop)
    _apply_image_transparency(image, transparency)

    return image

def _validate_image_path(image_path: Path):
    """Valide l'existence du fichier image"""
    if not image_path.exists():
        raise FileNotFoundError(f"Image non trouvée : {image_path}")

def _calculate_image_dimensions_with_aspect_ratio(image_path: Path, position):
    """Calcule les dimensions en maintenant le ratio d'aspect"""
    with Image.open(image_path) as img:
        width, height = img.size
        ratio = width / height

        if position is None:
            # Dimensions par défaut avec centrage
            if width > height:
                width = min(width, Inches(10).pt)
                height = width / ratio
            else:
                height = min(height, Inches(7).pt)
                width = height * ratio

            x = (Inches(10).pt - width) / 2
            y = (Inches(7.5).pt - height) / 2
            position = (x / Inches(1).pt, y / Inches(1).pt,
                      width / Inches(1).pt, height / Inches(1).pt)

    return position

def _setup_default_image_position():
    """Configure la position par défaut de l'image"""
    return (2, 2, 6, 4.5)  # x, y, width, height en pouces

def _add_image_to_slide(slide, image_path: Path, position):
    """Ajoute l'image à la slide avec la position calculée"""
    x, y, cx, cy = position
    return slide.shapes.add_picture(
        str(image_path),
        Inches(x), Inches(y),
        Inches(cx), Inches(cy)
    )

def _apply_image_cropping(image, crop):
    """Applique le recadrage à l'image si spécifié"""
    if crop:
        top, right, bottom, left = crop
        image.crop_top = top
        image.crop_right = right
        image.crop_bottom = bottom
        image.crop_left = left

def _apply_image_transparency(image, transparency):
    """Applique la transparence à l'image si spécifiée"""
    if transparency is not None:
        if not 0 <= transparency <= 1:
            raise ValueError("La transparence doit être comprise entre 0 et 1")
        image.transparency = int(transparency * 100)

def get_slide_info(slide) -> Dict[str, Any]:
    """Récupère les informations d'une slide.
    
    Args:
        slide: La slide à analyser
        
    Returns:
        Dictionnaire avec les informations de la slide
    """
    info = {
        "type": slide.slide_layout.name,
        "title": None,  # Initialisation à None
        "shapes": []
    }
    
    # Récupération du titre s'il existe et n'est pas vide
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        title_text = slide.shapes.title.text.strip()
        info["title"] = title_text if title_text else None
    
    for shape in slide.shapes:
        shape_info = {
            "type": shape.shape_type,
            "name": shape.name,
            "text": shape.text.strip() if hasattr(shape, "text") else None
        }
        info["shapes"].append(shape_info)
        
    return info 

__all__ = ['format_text']
