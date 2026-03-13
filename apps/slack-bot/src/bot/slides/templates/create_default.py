"""
Script to create the default presentation template.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def create_default_template():
    """Create the default presentation template."""
    prs = Presentation()
    
    # Title slide layout
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Title Slide"
    subtitle.text = "Subtitle"
    
    # Content slide layout
    content_slide = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Content Slide"
    content.text = "• Bullet point 1\n• Bullet point 2\n• Bullet point 3"
    
    # Chart slide layout
    chart_slide = prs.slide_layouts[5]
    slide = prs.slides.add_slide(chart_slide)
    title = slide.shapes.title
    title.text = "Chart Slide"
    
    # Image slide layout (blank)
    image_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(image_slide)
    # Add title shape since blank layout doesn't have one
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1.0)
    title = slide.shapes.add_textbox(left, top, width, height)
    title.text_frame.text = "Image Slide"
    
    # Save template
    template_dir = Path(__file__).parent
    prs.save(template_dir / "default.pptx")

if __name__ == "__main__":
    create_default_template() 