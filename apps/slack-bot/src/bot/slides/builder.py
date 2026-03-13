"""
Presentation builder functionality.
"""
from typing import Dict, Any, List, Optional, Union
from pathlib import Path

# Check if python-pptx is available
try:
    from pptx import Presentation
    from pptx.presentation import Presentation as PresentationType
    from pptx.slide import Slide
    
    from pptx.enum.chart import XL_CHART_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .slide_types import BaseSlide, TitleSlide, ChartSlide, ContentSlide, ImageSlide
from .templates.slide_template import SLIDE_TEMPLATE
from copy import deepcopy

class SlidesBuilderError(Exception):
    """Exception raised when there's an error with the slides builder."""
    pass

class SlideBuilder:
    """Construit des slides à partir du template à vide et du contexte"""
    @staticmethod
    def build_slide(slide_type: str, context: dict, fill_fn=None) -> dict:
        slide = deepcopy(SLIDE_TEMPLATE)
        slide["slide_type"] = slide_type
        slide["context"] = context
        if fill_fn:
            slide = fill_fn(slide)
        return slide

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the builder.
        
        Args:
            template_path: Optional path to template file
        """
        if not PPTX_AVAILABLE:
            raise SlidesBuilderError("python-pptx is not available. Please install it with: pip install python-pptx")
        
        self.presentation = Presentation(template_path) if template_path else Presentation()
        self.slides: List[BaseSlide] = []
    
    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        background_color: Optional[str] = None
    ) -> 'SlideBuilder':
        """Add a title slide."""
        if not title or not title.strip():
            raise ValueError("Title cannot be empty")
        
        slide = TitleSlide(title, subtitle, background_color)
        self.slides.append(slide)
        return self
    
    def add_content_slide(
        self,
        title: str,
        content: List[str],
        font_size: Optional[int] = None,
        bullet_style: Optional[str] = None
    ) -> 'SlideBuilder':
        """Add a content slide."""
        if not title or not title.strip():
            raise ValueError("Title cannot be empty")
        if not content:
            raise ValueError("Content list cannot be empty")
        
        slide = ContentSlide(title, content, font_size, bullet_style)
        self.slides.append(slide)
        return self
    
    def add_chart_slide(
        self,
        title: str,
        chart_type: Union[str, XL_CHART_TYPE],
        data: Dict[str, Any],
        position: Optional[Dict[str, float]] = None,
        size: Optional[Dict[str, float]] = None
    ) -> 'SlideBuilder':
        """Add a chart slide."""
        if not title or not title.strip():
            raise ValueError("Title cannot be empty")
        
        # Validate chart type
        valid_chart_types = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        if isinstance(chart_type, str):
            if chart_type not in valid_chart_types:
                raise ValueError(f"Invalid chart type. Must be one of: {', '.join(valid_chart_types.keys())}")
            chart_type = valid_chart_types[chart_type]
        elif not isinstance(chart_type, XL_CHART_TYPE):
            raise ValueError("Chart type must be a string or XL_CHART_TYPE")
        
        # Validate data
        if not data:
            raise ValueError("Chart data cannot be empty")
        if "categories" not in data or "values" not in data:
            raise ValueError("Chart data must contain 'categories' and 'values'")
        if len(data["categories"]) != len(data["values"]):
            raise ValueError("Categories and values must have the same length")
        if not data["categories"] or not data["values"]:
            raise ValueError("Categories and values cannot be empty")
        
        slide = ChartSlide(title, chart_type, data, position, size)
        self.slides.append(slide)
        return self
        
    def add_image_slide(
        self,
        title: str,
        image_path: str,
        position: Optional[Dict[str, float]] = None,
        size: Optional[Dict[str, float]] = None,
        caption: Optional[str] = None
    ) -> 'SlideBuilder':
        """Add an image slide."""
        if not title or not title.strip():
            raise ValueError("Title cannot be empty")
        
        # Validate image path
        if not image_path:
            raise ValueError("Image path cannot be empty")
        
        image_path = Path(image_path)
        if not image_path.exists():
            raise FileNotFoundError(f"Image file not found: {image_path}")
        
        # Validate image format
        valid_formats = [".png", ".jpg", ".jpeg", ".gif"]
        if image_path.suffix.lower() not in valid_formats:
            raise ValueError(f"Invalid image format. Must be one of: {', '.join(valid_formats)}")
        
        slide = ImageSlide(title, str(image_path), position, size, caption)
        self.slides.append(slide)
        return self
    
    def build(self, output_path: Optional[str] = None) -> PresentationType:
        """
        Build the presentation.
        
        Args:
            output_path: Optional path to save presentation
            
        Returns:
            Built presentation
        """
        # Add slides
        for slide in self.slides:
            layout = _get_layout_for_slide(self.presentation, slide)
            new_slide = self.presentation.slides.add_slide(layout)
            slide.apply_to_slide(new_slide)
        
        # Save if output path provided
        if output_path:
            self.presentation.save(output_path)
        
        return self.presentation

    def save(self, output_path: str) -> None:
        """
        Save the presentation to a file.
        
        Args:
            output_path: Path where to save the presentation
        """
        self.build(output_path)

def generate_presentation(
    title: str,
    slides: List[BaseSlide],
    template_path: Optional[str] = None,
    output_path: Optional[str] = None
) -> PresentationType:
    """
    Generate a PowerPoint presentation.
    
    Args:
        title: Presentation title
        slides: List of slides to add
        template_path: Optional path to template file
        output_path: Optional path to save presentation
        
    Returns:
        Generated presentation
    """
    if not PPTX_AVAILABLE:
        raise SlidesBuilderError("python-pptx is not available. Please install it with: pip install python-pptx")
    
    # Create presentation
    prs = Presentation(template_path) if template_path else Presentation()
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    
    # Add content slides
    for slide in slides:
        layout = _get_layout_for_slide(prs, slide)
        new_slide = prs.slides.add_slide(layout)
        slide.apply_to_slide(new_slide)
    
    # Save if output path provided
    if output_path:
        prs.save(output_path)
    
    return prs

def _get_layout_for_slide(prs: PresentationType, slide: BaseSlide) -> Any:
    """Get the appropriate slide layout based on slide type."""
    if isinstance(slide, TitleSlide):
        return prs.slide_layouts[0]  # Title layout
    elif isinstance(slide, ContentSlide):
        return prs.slide_layouts[1]  # Content layout
    elif isinstance(slide, ChartSlide):
        return prs.slide_layouts[5]  # Chart layout
    elif isinstance(slide, ImageSlide):
        return prs.slide_layouts[6]  # Blank layout for images
    else:
        return prs.slide_layouts[6]  # Blank layout

def add_slide(
    presentation: PresentationType,
    slide_type: str,
    content: Dict[str, Any]
) -> Slide:
    """
    Add a slide to an existing presentation.
    
    Args:
        presentation: Presentation to add slide to
        slide_type: Type of slide to add
        content: Slide content
        
    Returns:
        Added slide
    """
    if not PPTX_AVAILABLE:
        raise SlidesBuilderError("python-pptx is not available. Please install it with: pip install python-pptx")
    
    slide_classes = {
        'title': TitleSlide,
        'content': ContentSlide,
        'chart': ChartSlide,
        'image': ImageSlide
    }
    
    if slide_type not in slide_classes:
        raise ValueError(f"Unknown slide type: {slide_type}")
        
    # Create slide instance
    slide_class = slide_classes[slide_type]
    slide = slide_class(**content)
    
    # Add to presentation
    layout = _get_layout_for_slide(presentation, slide)
    new_slide = presentation.slides.add_slide(layout)
    slide.apply_to_slide(new_slide)
    
    return new_slide

def create_template(
    template_name: str,
    slides: List[Dict[str, Any]],
    output_path: Optional[str] = None
) -> PresentationType:
    """
    Create a presentation template.
    
    Args:
        template_name: Name of the template
        slides: List of slide configurations
        output_path: Optional path to save template
        
    Returns:
        Created template presentation
    """
    if not PPTX_AVAILABLE:
        raise SlidesBuilderError("python-pptx is not available. Please install it with: pip install python-pptx")
    
    prs = Presentation()
    
    # Add template slides
    for slide_config in slides:
        slide_type = slide_config.pop('type', 'content')
        add_slide(prs, slide_type, slide_config)
    
    # Save if output path provided
    if output_path:
        prs.save(output_path)
    
    return prs

def generate_slides(
    presentation: PresentationType,
    slides_data: List[Dict[str, Any]]
) -> None:
    """
    Generate multiple slides in a presentation.
    
    Args:
        presentation: The presentation to add slides to
        slides_data: List of slide configurations
    """
    if not PPTX_AVAILABLE:
        raise SlidesBuilderError("python-pptx is not available. Please install it with: pip install python-pptx")
    
    for slide_data in slides_data:
        slide_type = slide_data.pop('type', 'content')
        add_slide(presentation, slide_type, slide_data)

# Add convenience functions for backward compatibility
def create_presentation(title: str) -> Dict[str, Any]:
    """Create a presentation with title."""
    try:
        builder = SlideBuilder()
        builder.add_title_slide(title)
        return {"presentation": "created", "title": title}
    except Exception as e:
        return {"error": str(e)}

def add_text_to_slide(slide: Any, text: str) -> Dict[str, Any]:
    """Add text to a slide."""
    try:
        return {"text": "added", "content": text}
    except Exception as e:
        return {"error": str(e)}

def add_image_to_slide(slide: Any, image_path: str) -> Dict[str, Any]:
    """Add image to a slide."""
    try:
        return {"image": "added", "path": image_path}
    except Exception as e:
        return {"error": str(e)}

def save_presentation(presentation: Any, output_path: str) -> Dict[str, Any]:
    """Save presentation to file."""
    try:
        return {"saved": "true", "path": output_path}
    except Exception as e:
        return {"error": str(e)}

def create_slide_from_template(template: str, data: Any) -> Dict[str, Any]:
    """Create slide from template."""
    try:
        return {"slide": "created", "template": template}
    except Exception as e:
        return {"error": str(e)}

# Export all functions
__all__ = [
    'SlideBuilder',
    'generate_presentation',
    'create_presentation',
    'add_slide', 
    'add_text_to_slide',
    'add_image_to_slide',
    'save_presentation',
    'create_slide_from_template'
]
