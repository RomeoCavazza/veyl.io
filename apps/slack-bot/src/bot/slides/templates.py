"""
Template management functionality.
"""
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
import json
import logging

logger = logging.getLogger(__name__)

def get_template_path(template_name: str) -> Path:
    """
    Get the path to a template.
    
    Args:
        template_name: Name of the template
        
    Returns:
        Path to the template file
    """
    template_dir = Path(__file__).parent / "templates"
    template_path = template_dir / f"{template_name}.pptx"
    
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_name}")
    
    return template_path

def get_default_template() -> Dict[str, Any]:
    """
    Get the default canonical template.
    
    Returns:
        Default canonical template data
    """
    try:
        from .templates.canonical_templates import CANONICAL_TEMPLATES
        return CANONICAL_TEMPLATES.get('cover', {})
    except ImportError:
        # Fallback to file-based template
        template_dir = Path(__file__).parent / "templates"
        default_path = template_dir / "default.pptx"
        
        if not default_path.exists():
            raise FileNotFoundError("Default template not found")
        
        return {"path": str(default_path)}

def get_available_templates() -> List[str]:
    """
    Get list of available templates.
    
    Returns:
        List of template names
    """
    try:
        from .templates.canonical_templates import CANONICAL_TEMPLATES
        return list(CANONICAL_TEMPLATES.keys())
    except ImportError:
        # Fallback to file-based templates
        template_dir = Path(__file__).parent / "templates"
        if not template_dir.exists():
            return []
            
        return [f.stem for f in template_dir.glob("*.pptx")]

def get_template_metadata(template_name: str) -> Dict[str, Any]:
    """
    Get metadata for a template - refactorisé pour réduire la complexité

    Args:
        template_name: Name of the template

    Returns:
        Template metadata
    """
    try:
        # Essai 1: Template canonique
        metadata = _get_canonical_template_metadata(template_name)
        if metadata:
            return metadata

        # Essai 2: Template fichier
        metadata = _get_file_template_metadata(template_name)
        if metadata:
            return metadata

    except Exception as e:
        logger.error(f"Error getting template metadata: {e}")

    # Essai 3: Métadonnées par défaut
    return _get_default_metadata(template_name)

def _get_canonical_template_metadata(template_name: str) -> Optional[Dict[str, Any]]:
    """Récupère les métadonnées d'un template canonique"""
    try:
        from .templates.canonical_templates import CANONICAL_TEMPLATES
        template_data = CANONICAL_TEMPLATES.get(template_name, {})
        if template_data:
            return {
                "name": template_name,
                "description": f"Canonical template: {template_name}",
                "author": "REVOLVR AI",
                "version": "2.0",
                "type": "canonical",
                "data": template_data
            }
    except ImportError:
        pass
    return None

def _get_file_template_metadata(template_name: str) -> Optional[Dict[str, Any]]:
    """Récupère les métadonnées d'un template fichier"""
    template_dir = Path(__file__).parent / "templates"
    template_path = template_dir / f"{template_name}.pptx"

    if not template_path.exists():
        return None

    metadata_path = template_path.with_suffix('.json')
    return _read_metadata_file(metadata_path, template_name)

def _read_metadata_file(metadata_path: Path, template_name: str) -> Optional[Dict[str, Any]]:
    """Lit le fichier de métadonnées"""
    if not metadata_path.exists():
        return None

    try:
        with open(metadata_path) as f:
            return json.load(f)
    except Exception as e:
        logger.error(f"Error reading template metadata: {e}")
        return None

def _get_default_metadata(template_name: str) -> Dict[str, Any]:
    """Retourne les métadonnées par défaut"""
    return {
        "name": template_name,
        "description": f"Template {template_name}",
        "author": "System",
        "version": "1.0"
    }

def apply_theme(presentation: Presentation, theme_name: str) -> None:
    """
    Apply a theme to a presentation.
    
    Args:
        presentation: The presentation to apply the theme to
        theme_name: Name of the theme to apply
    """
    theme_path = get_template_path(theme_name)
    theme_prs = Presentation(theme_path)
    
    # Copy slide masters
    for master in theme_prs.slide_masters:
        # Copy master to presentation
        new_master = presentation.slide_masters.add_slide_master()
        
        # Copy master properties
        new_master.background = master.background
        new_master.follow_master_background = master.follow_master_background
        new_master.preserve_text = master.preserve_text
        
        # Copy layouts
        for layout in master.slide_layouts:
            new_layout = new_master.slide_layouts.add_slide_layout(layout.name)
            new_layout.background = layout.background
            new_layout.follow_master_background = layout.follow_master_background
            new_layout.preserve_text = layout.preserve_text
            
            # Copy placeholders
            for shape in layout.placeholders:
                new_shape = new_layout.placeholders.add_placeholder(
                    ph_type=shape.placeholder_format.type,
                    name=shape.name,
                    position=(shape.left, shape.top),
                    size=(shape.width, shape.height)
                )
                if shape.has_text_frame:
                    new_shape.text = shape.text

def get_layout(presentation: Presentation, layout_name: str) -> Any:
    """
    Get a slide layout by name - refactorisé pour réduire la complexité

    Args:
        presentation: The presentation to get the layout from
        layout_name: Name of the layout

    Returns:
        The slide layout
    """
    layout_name = layout_name.lower()

    # Essai 1: Recherche par nom exact
    layout = _find_layout_by_exact_name(presentation, layout_name)
    if layout:
        return layout

    # Essai 2: Recherche par index commun
    layout = _find_layout_by_index(presentation, layout_name)
    if layout:
        return layout

    # Essai 3: Recherche par nom partiel
    layout = _find_layout_by_partial_name(presentation, layout_name)
    if layout:
        return layout

    raise ValueError(f"Layout not found: {layout_name}")

def _get_common_layout_mapping() -> dict:
    """Retourne le mapping des layouts communs"""
    return {
        "title slide": 0,  # Title Slide
        "content slide": 1,  # Content with Caption
        "section header": 2,  # Section Header
        "two content": 3,  # Two Content
        "comparison": 4,  # Comparison
        "title only": 5,  # Title only
        "blank": 6,  # Blank
        "picture with caption": 7,  # Picture with Caption
        "chart": 8  # Chart
    }

def _find_layout_by_exact_name(presentation: Presentation, layout_name: str) -> Any:
    """Recherche un layout par nom exact"""
    for master in presentation.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.lower() == layout_name:
                return layout
    return None

def _find_layout_by_index(presentation: Presentation, layout_name: str) -> Any:
    """Recherche un layout par index commun"""
    common_layouts = _get_common_layout_mapping()

    if layout_name in common_layouts:
        try:
            return presentation.slide_layouts[common_layouts[layout_name]]
        except IndexError:
            pass

    return None

def _find_layout_by_partial_name(presentation: Presentation, layout_name: str) -> Any:
    """Recherche un layout par nom partiel"""
    for master in presentation.slide_masters:
        for layout in master.slide_layouts:
            if layout_name in layout.name.lower():
                return layout
    return None

__all__ = [
    'get_template_path',
    'get_default_template',
    'get_available_templates',
    'get_template_metadata',
    'apply_theme',
    'get_layout'
]
