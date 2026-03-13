"""PowerPoint Slide Builder for Revolver-AI-Bot
===========================================
Handles the generation of PowerPoint presentations with modular templates.
"""

import os
from pathlib import Path
from typing import Any, Dict

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from reco.models import DeckData

# Constants for slide layouts
LAYOUT_TITLE = 0  # Title slide layout
LAYOUT_CONTENT = 1  # Content slide layout

class SlideTemplate:
    """Base class for slide templates"""
    def __init__(self, presentation: PresentationType):
        self.presentation = presentation

    def apply(self, data: Dict[str, Any]) -> Slide:
        """Apply template with data"""
        raise NotImplementedError

class CoverSlide(SlideTemplate):
    """Template for cover slides"""
    def apply(self, data: Dict[str, Any]) -> Slide:
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[LAYOUT_TITLE]
        )
        title = slide.shapes.title
        if title:
            title.text = data.get("title", "")
        return slide

class InsightSlide(SlideTemplate):
    """Template for insight slides"""
    def apply(self, data: Dict[str, Any]) -> Slide:
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[LAYOUT_CONTENT]
        )
        title = slide.shapes.title
        if title:
            title.text = "Insight"
        
        # Add insight text in content placeholder
        content = slide.placeholders[1]
        content.text = data.get("insight", "")
        return slide

def build_slides(data: Dict[str, Any], output_path: str) -> Path:
    """Build a presentation from slide data"""
    prs = Presentation()
    
    # Map slide types to template classes
    templates = {
        "cover": CoverSlide,
        "insight": InsightSlide
    }
    
    # Create each slide
    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type")
        if slide_type in templates:
            template = templates[slide_type](prs)
            template.apply(slide_data)
    
    # Save presentation
    output = Path(output_path)
    prs.save(str(output))
    return output

# Aliases for backward compatibility
build_deck = build_slides
generate_deck = build_slides

def build_ppt(deck: DeckData, output_path: str):
    """Construit une présentation PPT à partir des données du deck - refactorisé"""
    prs = _setup_presentation()

    # Construction des slides section par section
    _add_brief_reminder_slide(prs, deck)
    _add_brand_overview_slide(prs, deck)
    _add_state_of_play_slides(prs, deck)
    _add_ideas_slides(prs, deck)
    _add_timeline_slide(prs, deck)
    _add_budget_slide(prs, deck)

    prs.save(output_path)

def _setup_presentation():
    """Configure la présentation avec le template approprié"""
    template_path = "pptx_generator/templates/base.pptx"
    return (
        Presentation(template_path) if os.path.exists(template_path) else Presentation()
    )

def _add_brief_reminder_slide(prs, deck: DeckData):
    """Ajoute la slide Brief Reminder"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Brief Reminder"

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    # Objectifs
    for obj in deck.brief_reminder.objectives:
        p = tf.add_paragraph()
        p.text = f"• {obj}"

    # Reformulation interne
    p = tf.add_paragraph()
    p.text = deck.brief_reminder.internal_reformulation

def _add_brand_overview_slide(prs, deck: DeckData):
    """Ajoute la slide Brand Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Brand Overview"

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for paragraph in deck.brand_overview.description_paragraphs:
        p = tf.add_paragraph()
        p.text = paragraph

def _add_state_of_play_slides(prs, deck: DeckData):
    """Ajoute les slides State of Play"""
    for idx, section in enumerate(deck.state_of_play, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"3.{idx} State of Play – {section.theme}"

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for ev in section.evidence:
            p = tf.add_paragraph()
            p.text = f"• {ev}"

def _add_ideas_slides(prs, deck: DeckData):
    """Ajoute les slides Ideas"""
    for idx, idea in enumerate(deck.ideas, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"4. Idea #{idx}: {idea.label}"

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for bullet in idea.bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"

# Helper pour créer des slides de contenu standardisées
def _create_content_slide(prs, title: str, content_items: List[str]) -> Slide:
    """Helper pour créer une slide de contenu standardisée."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item

    return slide

def _add_timeline_slide(prs, deck: DeckData):
    """Ajoute la slide Timeline avec helper optimisé."""
    timeline_items = [f"{m.deadline}: {m.label}" for m in deck.timeline]
    _create_content_slide(prs, "5. Timeline", timeline_items)

def _add_budget_slide(prs, deck: DeckData):
    """Ajoute la slide Budget avec helper optimisé."""
    budget_items = [f"{b.category}: €{b.estimate} ({b.comment})" for b in deck.budget]
    _create_content_slide(prs, "6. Budget", budget_items)
